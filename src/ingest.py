# src/ingest.py
import os

# ── PDF ──────────────────────────────────────────────────────────────────────
def extract_pages_from_pdf(path):
    try:
        from pypdf import PdfReader
    except ImportError:
        from PyPDF2 import PdfReader

    reader = PdfReader(path)
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append({"text": text.strip(), "source": os.path.basename(path), "page": i})
    return pages

# ── DOCX ─────────────────────────────────────────────────────────────────────
def extract_pages_from_docx(path):
    import docx
    doc = docx.Document(path)
    full_text = "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    if not full_text.strip():
        return []
    # treat whole doc as one "page"
    return [{"text": full_text.strip(), "source": os.path.basename(path), "page": 1}]

# ── PPTX ─────────────────────────────────────────────────────────────────────
def extract_pages_from_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    pages = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            pages.append({"text": "\n".join(texts), "source": os.path.basename(path), "page": i})
    return pages

# ── TXT ──────────────────────────────────────────────────────────────────────
def extract_pages_from_txt(path):
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    if not text.strip():
        return []
    return [{"text": text.strip(), "source": os.path.basename(path), "page": 1}]

# ── Router ────────────────────────────────────────────────────────────────────
def extract_pages(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        return extract_pages_from_pdf(path)
    elif ext == ".docx":
        return extract_pages_from_docx(path)
    elif ext in (".pptx", ".ppt"):
        return extract_pages_from_pptx(path)
    elif ext == ".txt":
        return extract_pages_from_txt(path)
    else:
        print(f"[ingest] Unsupported file type: {ext}")
        return []

# ── Chunker ───────────────────────────────────────────────────────────────────
def chunk_text(text, chunk_size=1000, overlap=200):
    chunks = []
    start = 0
    L = len(text)
    if L == 0:
        return []
    while start < L:
        end = min(start + chunk_size, L)
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start += chunk_size - overlap
    return chunks

# ── Ingest a folder ───────────────────────────────────────────────────────────
def ingest_folder(folder="data/books", chunk_size=1000, overlap=200):
    supported = (".pdf", ".docx", ".pptx", ".ppt", ".txt")
    all_chunks = []
    for fname in sorted(os.listdir(folder)):
        if not fname.lower().endswith(supported):
            continue
        path = os.path.join(folder, fname)
        print(f"[ingest] Reading {path}")
        pages = extract_pages(path)
        for p in pages:
            subchunks = chunk_text(p["text"], chunk_size=chunk_size, overlap=overlap)
            for n, sc in enumerate(subchunks, start=1):
                all_chunks.append({
                    "text": sc,
                    "source": p["source"],
                    "page": p["page"],
                    "chunk_id": f"{p['source']}_p{p['page']}_c{n}"
                })
    print(f"[ingest] Ingested {len(all_chunks)} chunks from {folder}")
    return all_chunks

# ── Ingest a single file ──────────────────────────────────────────────────────
def ingest_file(path, chunk_size=1000, overlap=200):
    print(f"[ingest] Reading single file: {path}")
    pages = extract_pages(path)
    all_chunks = []
    for p in pages:
        subchunks = chunk_text(p["text"], chunk_size=chunk_size, overlap=overlap)
        for n, sc in enumerate(subchunks, start=1):
            all_chunks.append({
                "text": sc,
                "source": p["source"],
                "page": p["page"],
                "chunk_id": f"{p['source']}_p{p['page']}_c{n}"
            })
    print(f"[ingest] Got {len(all_chunks)} chunks from {path}")
    return all_chunks

if __name__ == "__main__":
    chunks = ingest_folder()
    print("Sample chunk:", chunks[0] if chunks else "No chunks found")